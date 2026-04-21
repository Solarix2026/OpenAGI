# Copyright (c) 2026 ApeironAILab
# OpenAGI — Autonomous Intelligence System
# MIT License — https://github.com/ApeironAILab/OpenAGI

"""
html_ppt_builder.py — PPT generation pipeline:
1. NVIDIA generates slide content as JSON
2. Python renders each slide as styled HTML
3. Playwright screenshots each slide at 1920x1080
4. python-pptx assembles images into .pptx

This produces beautiful, design-quality slides vs. text-only python-pptx.
"""
import json, re, logging, tempfile, base64
from pathlib import Path
from core.llm_gateway import call_nvidia

log = logging.getLogger("HTMLPPTBuilder")
DECKS_DIR = Path("./workspace/decks")

# Slide HTML template — dark tech theme matching OpenAGI UI
SLIDE_TEMPLATE = """<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8">
<style>
* {{ margin:0; padding:0; box-sizing:border-box; }}
body {{
    width:1920px; height:1080px; overflow:hidden;
    background: {bg};
    font-family: 'Space Grotesk', 'Segoe UI', sans-serif;
    color: {text_color};
    display: flex; flex-direction: column; justify-content: center;
    padding: 80px 100px;
}}
.slide-number {{
    position:absolute; bottom:30px; right:50px;
    font-size:18px; opacity:0.3;
}}
.tag {{
    display:inline-block; background:rgba(59,130,246,0.15);
    border:1px solid rgba(59,130,246,0.3); color:#60a5fa;
    padding:4px 16px; border-radius:99px; font-size:20px;
    margin-bottom:24px; letter-spacing:0.05em; text-transform:uppercase;
}}
h1 {{ font-size:{title_size}px; font-weight:700; line-height:1.15; margin-bottom:24px; }}
h2 {{ font-size:36px; font-weight:600; margin-bottom:16px; opacity:0.85; }}
.body {{ font-size:28px; line-height:1.6; opacity:0.8; max-width:1400px; }}
.bullets {{ list-style:none; }}
.bullets li {{
    padding:12px 0; border-bottom:1px solid rgba(255,255,255,0.06);
    display:flex; align-items:flex-start; gap:16px;
}}
.bullets li::before {{ content:'→'; color:#3b82f6; flex-shrink:0; margin-top:2px; }}
.two-col {{ display:grid; grid-template-columns:1fr 1fr; gap:60px; }}
.accent {{ color:#3b82f6; }}
.gradient-text {{
    background:linear-gradient(135deg,#60a5fa,#a78bfa);
    -webkit-background-clip:text; -webkit-text-fill-color:transparent;
}}
.stat {{ font-size:80px; font-weight:800; color:#3b82f6; }}
.stat-label {{ font-size:22px; opacity:0.6; margin-top:8px; }}
</style>
<link href="https://fonts.googleapis.com/css2?family=Space+Grotesk:wght@400;600;700;800&display=swap" rel="stylesheet">
</head>
<body>
{content}
<div class="slide-number">{slide_num} / {total}</div>
</body>
</html>"""

THEMES = {
    "dark": {"bg": "linear-gradient(135deg, #020617 0%, #0f172a 100%)", "text_color": "#e2e8f0", "title_size": "72"},
    "dark_blue": {"bg": "linear-gradient(135deg, #0a1628 0%, #0f2744 100%)", "text_color": "#e2e8f0", "title_size": "72"},
    "midnight": {"bg": "linear-gradient(135deg, #000000 0%, #111827 100%)", "text_color": "#f1f5f9", "title_size": "72"},
    "light": {"bg": "#ffffff", "text_color": "#0f172a", "title_size": "64"},
}


def _render_slide_html(slide: dict, slide_num: int, total: int, theme: str = "dark") -> str:
    """Convert slide JSON to styled HTML."""
    t = THEMES.get(theme, THEMES["dark"])
    layout = slide.get("layout", "bullets")
    title = slide.get("title", "")
    subtitle = slide.get("subtitle", "")
    bullets = slide.get("bullets", [])
    body = slide.get("body", "")
    tag = slide.get("tag", "")
    stat = slide.get("stat", "")
    stat_label = slide.get("stat_label", "")

    if layout == "title":
        content = f"""<div>
        {f'<div class="tag">{tag}</div>' if tag else ''}
        <h1 class="gradient-text">{title}</h1>
        {f'<div class="body">{subtitle}</div>' if subtitle else ''}
        </div>"""
    elif layout == "stat":
        content = f"""<div>
        <h1>{title}</h1>
        <div class="stat">{stat}</div>
        <div class="stat-label">{stat_label}</div>
        </div>"""
    elif layout == "two_col":
        left = slide.get("left", [])
        right = slide.get("right", [])
        left_html = "".join(f"<li>{b}</li>" for b in left)
        right_html = "".join(f"<li>{b}</li>" for b in right)
        content = f"""<div>
        <h1>{title}</h1>
        <div class="two-col">
            <ul class="bullets">{left_html}</ul>
            <ul class="bullets">{right_html}</ul>
        </div>
        </div>"""
    else:  # bullets (default)
        bullets_html = "".join(f"<li>{b}</li>" for b in bullets)
        content = f"""<div>
        {f'<div class="tag">{tag}</div>' if tag else ''}
        <h1>{title}</h1>
        {f'<h2>{subtitle}</h2>' if subtitle else ''}
        <ul class="bullets body">{bullets_html}</ul>
        {f'<p class="body">{body}</p>' if body else ''}
        </div>"""

    return SLIDE_TEMPLATE.format(
        content=content, slide_num=slide_num, total=total, **t
    )


def generate_ppt(topic: str, style: str = "dark", num_slides: int = 10, save_path: str = None) -> dict:
    """Full pipeline: topic → JSON → HTML slides → screenshots → PPTX
    Returns: {"success", "pptx_path", "html_paths", "slides"}
    """
    DECKS_DIR.mkdir(parents=True, exist_ok=True)
    safe_name = re.sub(r'[^\w]', '_', topic.lower())[:25]

    # Step 1: Generate slide content with NVIDIA
    prompt = f"""Create a {num_slides}-slide professional presentation about: "{topic}"
Return ONLY valid JSON:
{{
    "title": "presentation title",
    "slides": [
        {{
            "slide": 1, "layout": "title",
            "tag": "optional category label",
            "title": "slide title", "subtitle": "subtitle or null",
            "bullets": ["point 1", "point 2"],
            "body": "paragraph text or null",
            "stat": "big number or null", "stat_label": "label for stat",
            "left": ["left col bullets"], "right": ["right col bullets"]
        }}
    ]
}}
Layout options: title (opening/section), bullets (default), two_col (comparison), stat (big number highlight)
First slide: layout=title. Last slide: summary + call to action.
Make content specific, data-driven, professional."""

    raw = call_nvidia([{"role": "user", "content": prompt}], max_tokens=2000)
    m = re.search(r'\{.*\}', raw, re.DOTALL)
    if not m:
        return {"success": False, "error": "Failed to generate slide content"}
    try:
        deck = json.loads(m.group(0))
    except json.JSONDecodeError:
        return {"success": False, "error": "Invalid JSON from NVIDIA"}

    slides = deck.get("slides", [])
    if not slides:
        return {"success": False, "error": "No slides generated"}

    # Step 2: Render HTML files
    html_dir = DECKS_DIR / f"{safe_name}_html"
    html_dir.mkdir(exist_ok=True)
    html_paths = []
    for i, slide in enumerate(slides):
        html_content = _render_slide_html(slide, i + 1, len(slides), style)
        html_path = html_dir / f"slide_{i+1:02d}.html"
        html_path.write_text(html_content, encoding="utf-8")
        html_paths.append(str(html_path))

    # Step 3: Screenshot with Playwright (if available)
    img_paths = []
    try:
        from playwright.sync_api import sync_playwright
        with sync_playwright() as pw:
            browser = pw.chromium.launch(headless=True)
            page = browser.new_page(viewport={"width": 1920, "height": 1080})
            for j, html_path in enumerate(html_paths):
                page.goto(f"file:///{html_path.replace(chr(92), '/')}")
                page.wait_for_load_state("networkidle")
                page.wait_for_timeout(500)  # Let fonts/gradients render
                img_path = html_dir / f"slide_{j+1:02d}.png"
                page.screenshot(path=str(img_path), full_page=False)
                img_paths.append(str(img_path))
                log.info(f"[PPT] Screenshot {j+1}/{len(html_paths)}")
            browser.close()
    except Exception as e:
        log.warning(f"Playwright screenshot failed: {e}. PPTX will be image-less.")

    # Step 4: Build PPTX
    pptx_path = None
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image as PILImage

        prs = Presentation()
        prs.slide_width = Inches(20)
        prs.slide_height = Inches(11.25)
        blank_layout = prs.slide_layouts[6]  # Blank

        for k, slide_data in enumerate(slides):
            prs_slide = prs.slides.add_slide(blank_layout)
            if k < len(img_paths) and img_paths:
                # Add screenshot as full-slide image
                pic = prs_slide.shapes.add_picture(
                    img_paths[k], Inches(0), Inches(0),
                    prs.slide_width, prs.slide_height
                )
            else:
                # Fallback: text-only slide
                txBox = prs_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(18), Inches(9))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.add_paragraph()
                p.text = slide_data.get("title", "")

        pptx_path = str(DECKS_DIR / f"{safe_name}.pptx")
        prs.save(pptx_path)
        log.info(f"[PPT] Saved: {pptx_path}")
    except ImportError as e:
        log.warning(f"python-pptx not available: {e}")

    return {
        "success": True,
        "topic": topic,
        "slides": len(slides),
        "pptx_path": pptx_path,
        "html_dir": str(html_dir),
        "html_paths": html_paths,
        "img_paths": img_paths,
        "note": "Open the HTML files for best visual quality. PPTX is for sharing."
    }


def register_html_ppt_tool(registry):
    """Register the generate_ppt tool."""
    def make_ppt(params: dict) -> dict:
        topic = params.get("topic", "")
        if not topic:
            return {"success": False, "error": "Provide topic"}
        style = params.get("style", "dark")
        slides = int(params.get("slides", 10))
        return generate_ppt(topic, style, slides)

    registry.register(
        "generate_ppt",
        make_ppt,
        "Generate a beautiful HTML-rendered PowerPoint presentation. Produces styled slides with gradients, then exports to PPTX.",
        {
            "topic": {"type": "string", "required": True},
            "style": {"type": "string", "default": "dark", "description": "dark, dark_blue, midnight, light"},
            "slides": {"type": "integer", "default": 10}
        },
        "generation"
    )
